"""
Nautilus Pitch Deck Generator
使用python-pptx生成完整的16页投资人路演PPT

安装依赖:
pip install python-pptx pillow

使用方法:
python generate_pitch_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0, 102, 255),      # #0066FF
    'dark_blue': RGBColor(0, 51, 102),          # #003366
    'gold': RGBColor(255, 184, 0),              # #FFB800
    'white': RGBColor(255, 255, 255),           # #FFFFFF
    'dark_gray': RGBColor(51, 51, 51),          # #333333
    'light_gray': RGBColor(240, 240, 240),      # #F0F0
}

def create_presentation():
    ""创建演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs

def add_title_slide(prs):
    """第1页: 封面"""
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)

    # 背景渐变（手动添加形状）
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark_blue']
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(3),
        Inches(12), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Nautilus"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.5),
        Inches(12), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Agent价值互联网平台"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = COLORS['gold']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 标语
    tagline_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(12), Inches(0.6)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "两篇顶级论文的延伸与精彩呈现"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(24)
    tagline_para.font.color.rgb = COLORS['white']
    tagline_para.alignment = PP_ALIGN.CENTER

    # 融资信息
    funding_box = slide.shapes.add_textbox(
        Inches(2), Inches(6.5),
        Inches(12), Inches(0.5)
    )
    funding_frame = funding_box.text_frame
    funding_frame.text = "种子轮融资 | $2M-$5M"
    funding_para = funding_frame.paragraphs[0]
    funding_para.font.size = Pt(20)
    funding_para.font.color.rgb = COLORS['gold']
    funding_para.alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    """第2页: 问题"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "AI Agent市场的四大痛点"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark_blue']

    # 4个问题 (2x2网格)
    problems = [
        ("💰 高成本", "20-30%手续费\n咨询公司极贵"),
        ("🐌 低效率", "人工处理慢\n缺乏自动化"),
        ("⚠️ 缺乏信任", "单点故障\n信息不透明"),
        ("🚫 无法学习", "Agent孤立\n经验无法积累")
    ]

    positions = [
        (2, 2.5), (9, 2.5),
        (2, 5.5), (9, 5.5)
    ]

    for (title, desc), (x, y) in zip(problems, positions):
        # 问题框
    box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
     Inches(x), Inches(y),
          Inches(5), Inches(2)
        )
        box.fill.solid()
      box.fill.fore_color.rgb = COLORS['light_gray']
        box.line.color.rgb = COLORS['primary_blue']
        box.line.width = Pt(2)

        # 标题
        text_frame = box.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark_blue']

        # 描述
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_before = Pt(10)

def add_solution_slide(prs):
    """第3页: 解决方案"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Nautilus = Epiplexity + DMAS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark_blue']

    # 流程图
    flow_items = [
        "两篇顶级论文融合",
        "↓",
        "Trinity Engine架构",
        "↓",
        "Agent价值互联网"
    ]

    y_pos = 2.5
    for item in flow_items:
        text_box = slide.shapes.add_textbox(
            Inches(4), Inches(y_pos),
          Inches(8), Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        p = text_frame.paragraphs[0]
        p.font.size = Pt(32) if item != "↓" else Pt(40)
        p.font.bold = True if item != "↓" else False
        p.font.color.rgb = COLORS['primary_blue']
     p.alignment = PP_ALIGN.CENTER
        y_pos += 0.8

    # 三大突破
    breakthroughs = [
        "✓ EvoMap机制 - 持续学习进化",
        "✓ 去中心化架构 - 无单点故障",
        "✓ 5%低费率 - vs 20-30%"
    ]

    y_pos = 6.5
    for item in breakthroughs:
     text_box = slide.shapes.add_textbox(
          Inches(3), Inches(y_pos),
        Inches(10), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = item
        p = text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['dark_gray']
        y_pos += 0.6

def add_academic_slide(prs):
    """第4页: 学术背书"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "基于两篇顶级AI研究论文"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark_blue']

  # 论文1
    paper1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2),
        Inches(5.5), Inches(4)
    )
    paper1_box.fill.solid()
    paper1_box.fill.fore_color.rgb = COLORS['light_gray']
    paper1_box.line.color.rgb = COLORS['primary_blue']
    paper1_box.line.width = Pt(3)

    text_frame = paper1_box.text_frame
    text_frame.text = "arXiv:2601.03220"
    p = text_frame.paragraphs[0]
  p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    p = text_frame.add_paragraph()
    p.text = "\nEpiplexity理论"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_gray']

    p = text_frame.add_paragraph()
    p.text = "\n知识涌现机制"
  p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']

    # 论文2
    paper2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(2),
        Inches(5.5), Inches(4)
    )
    paper2_box.fill.solid()
    paper2_box.fill.fore_color.rgb = COLORS['light_gray']
    paper2_box.line.color.rgb = COLORS['gold']
    paper2_box.line.width = Pt(3)

    text_frame = paper2_box.text_frame
    text_frame.text = "arXiv:2512.02410"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']

    p = text_frame.add_paragraph()
    p.text = "\nDMAS架构"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark_gray']

    p = text_frame.add_paragraph()
    p.text = "\n🏆 Best Paper Award"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['gold']

  # 底部文字
    bottom_box = slide.shapes.add_textbox(
        Inches(2), Inches(7),
      Inches(12), Inches(0.8)
    )
    bottom_frame = bottom_box.text_frame
    bottom_frame.text = "首个完整实现两篇论文的生产系统"
    bottom_para = bottom_frame.paragraphs[0]
    bottom_para.font.size = Pt(28)
  bottom_para.font.bold = True
    bottom_para.font.color.rgb = COLORS['primary_blue']
    bottom_para.alignment = PP_ALIGN.CENTER

def add_market_slide(prs):
    """第5页: 市场机会""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "$950B市场机会"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark_blue']

    # 市场数据
    market_data = [
        "AI服务市场: $150B → $500B (2026-2030)",
        "自由职业市场: $400B",
        "CAGR: 35%"
    ]

    y_pos = 2.5
    for data in market_data:
        text_box = slide.shapes.add_textbox(
            Inches(2), Inches(y_pos),
            Inches(12), Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = data
        p = text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['dark_gray']
        y_pos += 0.8

    # 目标客户
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(5),
        Inches(12), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "目标客户"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = COLORS['primary_blue']

    customers = [
        "👨‍💻 AI开发者: 5M+",
        "🏢 企业客户: 100K+",
        "🤖 AI Agent: 10M+"
    ]

    y_pos = 6
    for customer in customers:
        text_box = slide.shapes.add_textbox(
            Inches(3), Inches(y_pos),
         Inches(10), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = customer
      p = text_frame.paragraphs[0]
     p.font.size = Pt(24)
     p.font.color.rgb = COLORS['dark_gray']
        y_pos += 0.6

def add_remaing_slides(prs):
    """添加剩余11页（简化版本）"""

    slides_content = [
        {
       "title": "Trinity Engine - 三层架构",
          "content": [
                "Layer 1: Nexus Protocol",
                "• A2A通信 • Trust-Aware",
                "",
                "Layer 2: Orchestrator Engine",
       "• 智能调度 • 知识涌现",
       "",
             "Layer 3: Memory Chain",
                "• 多层记忆 • 区块链"
            ]
        },
        {
            "title": "多元化收入模式",
            "content": [
                "订阅收入 (60%)",
                "Free | Pro $29/月 | Enterprise",
                "",
              "平台交易费 (35%)",
          "5% vs 竞争对手20-30%",
                "",
                "企业定制 (5%)",
             "",
         "CAC: $10 | LTV: $500 | LTV/CAC: 50"
            ]
        },
        {
            "title": "五大核心优势",
            "content": [
                "学术背书: 2篇顶级论文",
                "费率: 5% vs 20-30%",
            "效率: AI自动化10倍",
                "信任: 去中心化",
             "学习: EvoMap持续进化",
            "",
                "护城河: 技术专利 | 网络效应 | 数据优势"
            ]
        },
      {
            "title": "已完成的里程碑",
            "content": [
           "✓ MVP完成",
            "✓ 系统稳定运行",
                "✓ 测试覆盖率70%",
                "✓ 技术白皮书发布",
                "",
       "早期用户",
           "• Beta测试用户",
                "• 开发者反馈",
          "• 企业试用"
       ]
        },
     {
            "title": "3年增长计划",
          "content": [
              "Year 1: 30K用户, $250K收入",
             "Year 2: 100K用户, $1.52M收入",
                "Year 3: 500K用户, $10.2M收入",
                "",
                "盈利时间: Year 3 Q2",
                "毛利率: 70%+",
                "LTV/CAC: 50"
            ]
        },
        {
            "title": "三阶段增长策略",
            "content": [
              "Phase 1: 种子用户 (M1-3)",
                "• Product Hunt • 开发者社区",
              "",
                "Phase 2: 增长期 (M4-12)",
              "• 付费广告 • 合作伙伴",
          "",
         "Phase 3: 规模化 (Y2-3)",
                "• 国际化 • 企业销售"
            ]
        },
        {
            "title": "经验丰富的核心团队",
       "content": [
        "技术团队",
                "• AI/ML专家",
                "• 区块链工程师",
                "• 分布式系统架构师",
           "",
                "商业团队",
                "• CEO: 10年+创业经验",
             "• COO: 运营管理专家",
        "• CMO: 增长营销专家",
                "",
             "学术顾问: 论文作者团队"
            ]
        },
        {
            "title": "融资$2M-$5M",
            "content": [
                "产品开发 (40%) $0.8M-$2M",
         "市场营销 (30%) $0.6M-$1.5M",
         "团队扩张 (20%) $0.4M-$1M",
                "运营资金 (10%) $0.2M-$0.5M",
                "",
                "估值: $10M-$20M (pre-money)",
                "稀释: 10-20%",
              "",
         "里程碑: 6M→10K用户 | 12M→盈利"
            ]
        },
        {
            "title": "成为Agent价值互联网基础设施",
            "content": [
                "短期 (1-2年)",
        "AI Agent领域领导者 | 10万+用户",
                "",
                "中期 (3-5年)",
              "全球最大Agent市场 | 100万+用户",
                "年收入$100M+",
                "",
                "长期 (5-10年)",
        "Internet of Agents实现",
                "推动AGI发展"
            ]
        },
        {
         "title": "让我们一起改变AI的未来",
            "content": [
             "网站: nautilus.social",
                "邮箱: invest@nautilus.social",
                "",
                "Twitter | LinkedIn | GitHub",
                "",
       "加入我们，共同打造",
                "Agent价值互联网！"
            ]
        },
        {
            "title": "补充材料",
            "content": [
                "可选内容:",
              "",
                "• 详细财务模型",
            "• 技术架构深度解析",
                "• 市场研究数据",
                "• 客户案例研究",
             "• 团队完整履历",
                "• 竞争对手详细分析",
             "",
                "完整材料请访问:",
                "nautilus.social/investor"
        ]
     }
    ]

    for slide_data in slides_content:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

      # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
         Inches(14), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data["title"]
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['dark_blue']

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(2), Inches(2),
        Inches(12), Inches(6)
        )
        content_frame = content_box.text_frame

        for i, line in enumerate(slide_data["content"]):
      if i == 0:
              content_frame.text = line
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
        p.text = line

            p.font.size = Pt(24) if line and not line.startswith("•") else Pt(20)
            p.font.color.rgb = COLORS['dark_gray']
            if line.startswith("✓") or line.startswith("•"):
             p.level = 1

def main():
  """主函数"""
    print("开始生成Nautilus Pitch Deck...")

    # 创建演示文稿
    prs = create_presentation()

    # 添加所有页面
    print("添加第1页: 封面...")
    add_title_slide(prs)

    print("添加第2页: 问题...")
    add_problem_slide(prs)

    print("添加第3页: 解决方案...")
    add_solution_slide(prs)

    print("添加第4页: 学术背书...")
    add_academic_slide(prs)

    print("添加第5页: 市场机会...")
    add_market_slide(prs)

    print("添加第6-16页...")
    add_remaining_slides(prs)

    # 保存文件
    filename = 'Nautilus_Pitch_Deck_CN.pptx'
    prs.save(filename)
    print(f"\n✅ PPT生成完成: {filename}")
    print(f"📄 总页数: {len(prs.slides)}")
    print("\n下一步:")
    print("1. 打开PPT文件")
    print("2. 添加Logo和品牌元素")
    print("3. 调整配色和字体")
    print("4. 添加动画效果")
    print("5. 导出为PDF")

if __name__ == "__main__":
    main()
