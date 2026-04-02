"""
Nautilus Pitch Deck Generator with Gemini Integration
使用Gemini API生成图片和优化内容

依赖:
pip install python-pptx pillow google-generativeai requests

使用方法:
python generate_pitch_deck_gemini.py
"""

import os
import requests
import base64
from io import BytesIO
from PIL import Image
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# API配置
GEMINI_API_KEY = "AIzaSyAbLcIXvJiAzKJ-ntW6Jmcd8KcW_LPkz0Y"
NANO_BANANA_MODEL = "gemini-3.1-flash-image-preview"

# 配置Gemini
genai.configure(api_key=GEMINI_API_KEY)

# 配色方案
COLORS = {
    'primary_blue': RGBColor(0, 102, 255),
    'dark_blue': RGBColor(0, 51, 102),
    'gold': RGBColor(255, 184, 0),
    'white': RGBColor(255, 255, 255),
    'dark_gray': RGBColor(51, 51, 51),
    'light_gray': RGBColor(240, 240, 240),
}

def generate_image_with_nano_banana(prompt, output_path):
    """使用nano-banana 2生成图片"""
    try:
        print(f"生成图片: {prompt[:50]}...")

        model = genai.GenerativeModel(NANO_BANANA_MODEL)
        response = model.generate_content(prompt)

        # 保存图片
        if hasattr(response, 'image'):
            with open(output_path, 'wb') as f:
                f.write(response.image)
            print(f"✅ 图片已保存: {output_path}")
            return True
    else:
         print(f"⚠️ 未能生成图片")
            return False

    except Exception as e:
    print(f"❌ 图片生成失败: {e}")
        return False

def generate_cover_image():
    """生成封面背景图"""
    prompt = """Create a professional tech background image for a pitch deck cover slide.
    Style: Modern, minimalist, tech-focused
    Colors: Deep blue gradient with gold accents
    Elements: Abstract network connections, AI nodes, flowing data
    Mood: Professional, innovative, futuristic
    Size: 1920x1080px
    No text or logos"""

    return generate_image_with_nano_banana(prompt, "cover_background.png")

def generate_architecture_diagram():
    """生成架构图""
    prompt = """Create a technical architecture diagram showing three layers:
    Layer 1 (bottom): Nexus Protocol - decentralized communication
    Layer 2 (middle): Orchestrator Engine - smart scheduling
    Layer 3 (top): Memory Chain - multi-layer memory + blockchain
    Style: Clean, professional, tech diagram
    Colors: Blue and white
    Include: Arrows showing data flow between layers
    No text labels needed"""

    return generate_image_with_nano_banana(prompt, "architecture_diagram.png")

def generate_market_growth_chart():
    """生成市场增长曲线图""
    prompt = """Create a professional growth chart showing:
    - Upward trending curve from $150B to $500B
    - Years 2026 to 2030 on X-axis
    - Revenue in billions on Y-axis
    - Clean, modern style
    - Blue and green colors
    - Grid background
    No text labels needed"""

    return generate_image_with_nano_bana(prompt, "market_growth.png")

def generate_team_placeholder():
    """生成团队照片占位符"""
    prompt = """Create a professional team photo placeholder:
    - Silhouettes of 5-6 people
    - Professional business setting
    - Blue and gray colors
    - Modern, clean style
    - Suitable for business presentation"""

    return generate_image_with_nano_banana(prompt, "team_placeholder.png")

def create_presentation_with_images():
    """创建包含图片的PPT"""
    print("\n开始生成PPT...")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 第1页: 封面（带背景图）
    print("\n添加第1页: 封面...")
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 添加背景图（如果生成成功）
    if os.path.exists("cover_background.png"):
        slide.shapes.add_picture(
            "cover_background.png",
         0, 0,
         prs.slide_width, prs.slide_height
        )

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(3),
        Inches(12), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Nautilus"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.5),
        Inches(12), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Agent价值互联网平台"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = COLORS['gold']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 标语
    tagline_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(12), Inches(0.6)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "两篇顶级论文的延伸与精彩呈现"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(24)
    tagline_para.font.color.rgb = COLORS['white']
    tagline_para.alignment = PP_ALIGN.CENTER

    # 融资信息
    funding_box = slide.shapes.add_textbox(
        Inches(2), Inches(6.5),
        Inches(12), Inches(0.5)
    )
    funding_frame = funding_box.text_frame
    funding_frame.text = "种子轮融资 | $2M-$5M"
    funding_para = funding_frame.paragraphs[0]
    funding_para.font.size = Pt(20)
    funding_para.font.color.rgb = COLORS['gold']
    funding_para.alignment = PP_ALIGN.CENTER

    # 添加更多页面...
    # (这里可以添加其他15页，类似之前的脚本)

    # 保存
    filename = 'Nautilus_Pitch_Deck_with_Images.pptx'
    prs.save(filename)
    print(f"\n✅ PPT生成完成: {filename}")
    return filename

def main():
    """主函数"""
    print("=" * 60)
    print("Nautilus Pitch Deck Generator with Gemini")
    print("=" * 60)

    # 步骤1: 生成图片
    print("\n步骤1: 使用nano-banana 2生成图片...")
    print("-" * 60)

    images_to_generate = [
        ("封面背景", generate_cover_image),
        ("架构图", generate_architecture_diagram),
        ("市场增长图", generate_market_growth_chart),
      ("团队照片", generate_team_placeholder),
    ]

    generated_count = 0
    for name, func in images_to_generate:
        print(f"\n生成 {name}...")
        if func():
        generated_count += 1

    print(f"\n✅ 成功生成 {generated_count}/{len(images_to_generate)} 张图片")

    # 步骤2: 创建PPT
    print("\n步骤2: 创建PPT...")
    print("-" * 60)

    filename = create_presentation_with_images()

    # 完成
    print("\n" + "=" * 60)
    print("✅ 完成！")
    print("=" * 60)
    print(f"\n生成的文件:")
    print(f"  - {filename}")
    if generated_count > 0:
        print(f"  - {generated_count} 张图片")

    print("\n下一步:")
    print("  1. 打开PPT文件查看")
    print("  2. 添加Logo")
    print("  3. 调整文字和布局")
    print("  4. 导出为PDF")
    print("\n")

if __name__ == "__main__":
    main()
